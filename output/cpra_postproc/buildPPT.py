#####################################################
# AUTHORS: Matthew V Bilskie, PhD
#          Louisiana State University
#
#          Jason Fleming, PhD
#          Seahorse Coastal Consulting
# COPYRIGHT 2018
#
#####################################################
#
# TO RUN IN A PYTHON SHELL USING CMD ARGS
#import sys
#sys.argv = ['./buildPPT.py','FigureGenFilename']
#execfile('./buildPPT.py')
#
#####################################################

import sys
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE


# Get command line argument
fname = sys.argv[1]

# Read run.properties and make a property dictionary
runProp = dict()
f = open('run.properties','r')
for line in f:
    fields = line.split(':',1)
    try:
        runProp[fields[0].strip()] = fields[1].strip()
    except IndexError:
        continue
f.close()

# Convert advisoryTime to python datetime object
advisory_dt = datetime.strptime(runProp['time.forecast.valid.cdt'],'%Y%m%d%H%M%S')
advisory_dt_long = datetime.strftime(advisory_dt,'%b-%d-%Y %H:%M')

scenario = runProp['asgs.enstorm']
if scenario == 'nhcConsensus':
    scenario = 'NHC Track'

prs = Presentation('LSU_template.pptx')

slide_layout = prs.slide_layouts[5]
slide_layout_hydro = prs.slide_layouts[6]

numSlides = 1

# Create a title slide
title_slide_layout = prs.slide_layouts[0]
#for shape in title_slide_layout.placeholders:
#    print('%d %s' % (shape.placeholder_format.idx, shape.name))
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = runProp['storm class'] + ' ' + runProp['stormname'] + ', ' + scenario + ' Scenario'
subtitle.text = "Advisory " + runProp['advisory'] + " Issued on " + advisory_dt_long + " CDT"
statement = 'For Official Use Only. Not For Release. \rModel results were produced by the ADCIRC Surge Guidance System (ASGS) and are based on the National Hurricane Center (NHC) forecast track. \rADCIRC-developed hydrographs are an operational planning tool for emergency-response personnel and are not a replacement for National Weather Service (NWS) forecasts.'
fouo = slide.placeholders[10]
fouo.text = statement
numSlides = numSlides + 1

# Set slide layout
left = Inches(1.94)
top = Inches(1.06)

img_path = fname
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = 'NHC Advisory ' + runProp['advisory'] + ' ' + scenario + ' Scenario'
subtitle.text = "Simulated peak water levels (ft, NAVD88)"
pic = slide.shapes.add_picture(img_path,left,top)
fouo = slide.placeholders[13]
fouo.text = statement
snum = slide.placeholders[14]
snum.text = str(numSlides)
numSlides = numSlides + 1
#for shape in slide.placeholders:
#    print('%d %s' % (shape.placeholder_format.idx, shape.name))

#left = Inches(0.42)
#top = Inches(1.15)
left = Inches(0.75)
top = Inches(0.81)
iwidth = Inches(11.84)
iheight = Inches(5.69)

fnames = ['WSE_17StCanal_USACE85625.png','WSE_IHNC01_USACE76065.png','WSE_IHNC02_USACE76030.png',
        'WSE_LPV144_USACE76010.png','WSE_LPV149_USACE85760.png','WSE_NOV13_USACE01440.png',
        'WSE_NOV14_USACE01440.png','WSE_WBV09a_USACE82770.png','WSE_WBV09b_USACE82762.png',
        'WSE_WBV162_USACE82742.png','WSE_WBV7274_USACE82715.png','WSE_WBV90_USACE76265.png',
        'WSE_LakefrontAirport_USACE85670.png','WSE_Mandeville_USACE85575.png',
        'WSE_Rigolets_USACE85700.png','WSE_Lafitte_USACE82875.png']

# Station names correspond to the order of fnames
#staName = ['17th St. Outfall Canal','Seabrook Complex (IHNC-01)','IHNC Surge Barrier (IHNC-02)',
#        'Bayou Dupre Sector Gate (LPV-144)','Caernarvon Canal Sector Gate (LPV-149)',
#        'Empire Floodgate (NOV-13)','Empire Lock (NOV-14)','Oakville Sluice Gate (WBV-09a)',
#        'Hero Canal stop-log gage (WBV-09b)','Bayou Segnetee closure (WBV-16.2)',
#        'Western Tie-In features (WBV-74-72)','West Closure Complex (WBV-90)',
#        'Lakefront Airport','Mandeville','Rigolets','Lafitte']
staName = ['Outfall 17th St London Ave Orleans Ave, LA (17StCanal, CPRA) (85625, USACE)',
           'Seabrook Complex - Flood Side, LA (IHNC01, CPRA) (76065, USACE)',
           'IHNC Surge Barrier East - Flood Side, LA (IHNC02, CPRA) (76030, USACE)',
           'Bayou Dupre Sector Gate - East/Flood Side, LA (LPV144, CPRA)\n(76010, USACE)',
           'Caernarvon Canal Sector Gate - South/Flood Side, LA (LPV149, CPRA)\n(85760, USACE)',
           'Mississippi River at Empire Floodgate, LA (NOV13, CPRA) (01440, USACE)',
           'Mississippi River at Empire Lock, LA (NOV14, CPRA) (01440, USACE)',
           'Oakville Sluice Gate - Flood Side/South (WBV09a, CPRA) (82770 USACE)',
           'Hero Canal Stop-Log Gate - Flood Side/West, LA (WBV09b, CPRA)\n(82762, USACE)',
           'Bayou Segnette Closure - Flood Side, LA (WBV-16.2, CPRA) (82742, USACE)',
           'Bayou Verret / W. Tie-In Sector Gate Flood Side, LA (WBV-72/74, CPRA)\n(82715, USACE)',
           'GIWW at West Closure Complex - Flood Side, LA (WBV90, CPRA)\n(76265, USACE)',
           'Lake Pontchartrain at Lakefront Airport, LA (LakefrontAirport, CPRA)\n(85670, USACE)',
           'Lake Pontchartrain at Mandeville, LA (Mandeville, CPRA) (85575, USACE)',
           'Rigolets near Lake Pontchartrain, LA (Rigolets, CPRA) (85700, USACE)',
           'Barataria Waterway at Lafitte, LA (Lafitte, CPRA) (82875, USACE)']

i = 0
for image in fnames:
    try: 
        slide = prs.slides.add_slide(slide_layout_hydro)
        title = slide.shapes.title
        title.text = staName[i]
        pic = slide.shapes.add_picture(image,left,top,width=iwidth,height=iheight)
        fouo = slide.placeholders[13]
        fouo.text = statement
        snum = slide.placeholders[14]
        snum.text = str(numSlides)
        numSlides = numSlides + 1
        i = i + 1
    except:
        print("ERROR: buildPPT.py: Could not find " + image + ".")

# Loop through slides
#slides = prs.slides
#for slide in slides:
        #print('slide number %s' % str(slides.index(slide)+1))

pptFile = runProp['stormname'] + "_Adv" + runProp['advisory'] + "_" + scenario + "_" + runProp['forecastValidStart'] + ".pptx"
prs.save(pptFile)
pFile = open('pptFile.temp','w')
pFile.write(pptFile)
pFile.close()
